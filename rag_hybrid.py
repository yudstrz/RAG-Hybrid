# ==========================
# 🌟 RAG SYSTEM - GEMINI (Fixed & Robust Version)
# ==========================
import os
import requests
import tempfile
import uuid
import logging
import json
from datetime import datetime
from typing import List, Dict, Any, Optional
import streamlit as st

# Core dependencies
import numpy as np

# Web search & scraping
try:
    from duckduckgo_search import DDGS
except ImportError:
    DDGS = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

# Document loaders
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx2txt
except ImportError:
    docx2txt = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Embeddings and vector store
try:
    from sentence_transformers import SentenceTransformer
except ImportError:
    SentenceTransformer = None

try:
    import chromadb
    from chromadb.config import Settings
except ImportError:
    chromadb = None

# LangChain for text splitting
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    RecursiveCharacterTextSplitter = None

# ==========================
# 📝 Logging Configuration
# ==========================
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ==========================
# 🎨 Streamlit Page Configuration
# ==========================
st.set_page_config(
    page_title="🤖 Gemini RAG Assistant",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ==========================
# ✅ Dependency Check
# ==========================
def check_dependencies():
    """Check if all required dependencies are installed"""
    missing = []
    
    if SentenceTransformer is None:
        missing.append("sentence-transformers")
    if chromadb is None:
        missing.append("chromadb")
    if RecursiveCharacterTextSplitter is None:
        missing.append("langchain-text-splitters")
    if DDGS is None:
        missing.append("duckduckgo-search")
    if BeautifulSoup is None:
        missing.append("beautifulsoup4 lxml")
    if pypdf is None:
        missing.append("pypdf")
    if docx2txt is None:
        missing.append("docx2txt")
    if Presentation is None:
        missing.append("python-pptx")
    
    if missing:
        st.error(f"❌ Missing dependencies: {', '.join(missing)}")
        st.code(f"pip install {' '.join(missing)}")
        return False
    return True

# ==========================
# 🔑 API Key Management (SECURE)
# ==========================
def get_api_key() -> Optional[str]:
    """Get API key with priority: secrets > env > user input"""
    try:
        # Priority 1: Streamlit secrets
        if hasattr(st, 'secrets') and "GEMINI_API_KEY" in st.secrets:
            return st.secrets["GEMINI_API_KEY"]
    except Exception as e:
        logger.warning(f"Could not access secrets: {e}")
    
    # Priority 2: Environment variable
    env_key = os.getenv("GEMINI_API_KEY")
    if env_key:
        return env_key
    
    # Priority 3: User input (fallback)
    if "api_key_input" not in st.session_state:
        st.session_state.api_key_input = ""
    
    api_key = st.sidebar.text_input(
        "🔑 Masukkan Gemini API Key:",
        type="password",
        value=st.session_state.api_key_input,
        help="Dapatkan API key gratis di: https://ai.google.dev/"
    )
    
    if api_key:
        st.session_state.api_key_input = api_key
        return api_key
    
    return None

# ==========================
# 🧩 Cached Resources
# ==========================
@st.cache_resource
def load_embedding_model():
    """Load embedding model with error handling"""
    if SentenceTransformer is None:
        return None
    
    try:
        logger.info("Loading embedding model...")
        model = SentenceTransformer("all-MiniLM-L6-v2")
        logger.info("Embedding model loaded successfully")
        return model
    except Exception as e:
        logger.error(f"Failed to load embedding model: {e}")
        st.error(f"❌ Gagal load embedding model: {e}")
        return None

@st.cache_resource
def load_chroma_client():
    """Initialize ChromaDB client"""
    if chromadb is None:
        return None
    
    try:
        logger.info("Initializing ChromaDB...")
        client = chromadb.Client(Settings(
            allow_reset=True,
            anonymized_telemetry=False
        ))
        logger.info("ChromaDB initialized successfully")
        return client
    except Exception as e:
        logger.error(f"Failed to initialize ChromaDB: {e}")
        st.error(f"❌ Gagal initialize ChromaDB: {e}")
        return None

def get_text_splitter():
    """Get text splitter with fallback"""
    if RecursiveCharacterTextSplitter:
        return RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=150,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
    else:
        # Simple fallback splitter
        class SimpleSplitter:
            def split_text(self, text, chunk_size=1000):
                chunks = []
                for i in range(0, len(text), chunk_size):
                    chunks.append(text[i:i+chunk_size])
                return chunks
        return SimpleSplitter()

# ==========================
# 📊 Session State Initialization
# ==========================
def init_session_state():
    """Initialize all session state variables"""
    defaults = {
        "initialized": False,
        "vector_store": None,
        "chat_history": [],
        "processing_stats": {},
        "last_query_time": None,
        "query_count": 0,
        "api_key_input": ""
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

# ==========================
# 🛡️ Rate Limiting
# ==========================
def check_rate_limit(max_queries_per_minute: int = 10) -> bool:
    """Simple rate limiting based on query count"""
    current_time = datetime.now()
    
    if st.session_state.last_query_time:
        time_diff = (current_time - st.session_state.last_query_time).total_seconds()
        
        # Reset counter after 1 minute
        if time_diff > 60:
            st.session_state.query_count = 0
            st.session_state.last_query_time = current_time
        
        # Check limit
        if st.session_state.query_count >= max_queries_per_minute:
            return False
    
    st.session_state.query_count += 1
    st.session_state.last_query_time = current_time
    return True

# ==========================
# 📄 Document Loading Functions
# ==========================
def load_and_split_documents(paths: List[str], text_splitter) -> List[Dict[str, Any]]:
    """Load and split documents with improved error handling"""
    docs = []
    stats = {"success": 0, "failed": 0, "total_chunks": 0}
    
    for p in paths:
        try:
            ext = p.lower().split(".")[-1]
            text_content = ""
            source_name = os.path.basename(p)
            
            # Load based on file type
            if ext == "pdf" and pypdf:
                try:
                    reader = pypdf.PdfReader(p)
                    text_content = "\n".join([
                        page.extract_text() or "" for page in reader.pages
                    ])
                except Exception as e:
                    logger.error(f"PDF read error for {source_name}: {e}")
                    stats["failed"] += 1
                    continue
                    
            elif ext in ["doc", "docx"] and docx2txt:
                try:
                    text_content = docx2txt.process(p)
                except Exception as e:
                    logger.error(f"DOCX read error for {source_name}: {e}")
                    stats["failed"] += 1
                    continue
                    
            elif ext in ["ppt", "pptx"] and Presentation:
                try:
                    prs = Presentation(p)
                    text_content = "\n".join([
                        shape.text for slide in prs.slides 
                        for shape in slide.shapes 
                        if hasattr(shape, "text")
                    ])
                except Exception as e:
                    logger.error(f"PPTX read error for {source_name}: {e}")
                    stats["failed"] += 1
                    continue
                    
            elif ext in ["txt", "md"]:
                try:
                    with open(p, "r", encoding="utf-8", errors="ignore") as f:
                        text_content = f.read()
                except Exception as e:
                    logger.error(f"TXT read error for {source_name}: {e}")
                    stats["failed"] += 1
                    continue
            else:
                logger.warning(f"Unsupported file type: {ext}")
                stats["failed"] += 1
                continue
            
            # Split text to chunks
            if text_content.strip():
                try:
                    chunks = text_splitter.split_text(text_content)
                    for i, chunk in enumerate(chunks):
                        docs.append({
                            "content": chunk,
                            "metadata": {
                                "source": source_name,
                                "chunk": i + 1,
                                "total_chunks": len(chunks),
                                "type": "document"
                            }
                        })
                    
                    stats["success"] += 1
                    stats["total_chunks"] += len(chunks)
                    logger.info(f"Loaded {source_name}: {len(chunks)} chunks")
                except Exception as e:
                    logger.error(f"Text split error for {source_name}: {e}")
                    stats["failed"] += 1
            else:
                logger.warning(f"Empty content in {source_name}")
                stats["failed"] += 1
                
        except Exception as e:
            logger.error(f"Failed to load {os.path.basename(p)}: {e}")
            stats["failed"] += 1
    
    st.session_state.processing_stats["documents"] = stats
    return docs

# ==========================
# 🌐 Web Search Functions
# ==========================
def search_and_scrape_web(
    query: str,
    text_splitter,
    max_results: int = 3,
    scrape_full_page: bool = True
) -> List[Dict[str, Any]]:
    """Search and scrape web with improved error handling"""
    results = []
    stats = {"success": 0, "failed": 0, "total_chunks": 0}
    
    if not DDGS:
        st.error("❌ duckduckgo-search tidak terinstal")
        return results
    
    try:
        # Search with DuckDuckGo
        with DDGS() as ddgs:
            ddgs_results = list(ddgs.text(query, max_results=max_results))
        
        logger.info(f"Found {len(ddgs_results)} search results")
        
        for idx, r in enumerate(ddgs_results, 1):
            try:
                content = r.get("body", "")
                title = r.get("title", "Untitled")
                url = r.get("href", "")
                
                # Scrape full page if requested
                if scrape_full_page and url and BeautifulSoup:
                    try:
                        with st.spinner(f"🕷️ Scraping {idx}/{len(ddgs_results)}: {title[:50]}..."):
                            response = requests.get(
                                url,
                                timeout=10,
                                headers={"User-Agent": "Mozilla/5.0"}
                            )
                            response.raise_for_status()
                            
                            soup = BeautifulSoup(response.content, "lxml")
                            
                            # Remove script and style tags
                            for tag in soup(["script", "style", "nav", "footer", "header"]):
                                tag.decompose()
                            
                            # Extract paragraphs
                            paragraphs = soup.find_all(['p', 'article', 'section'])
                            page_content = "\n".join([
                                p.get_text(strip=True) for p in paragraphs
                            ])
                            
                            if page_content.strip() and len(page_content) > len(content):
                                content = page_content
                                logger.info(f"Scraped {len(page_content)} chars from {url}")
                            
                            stats["success"] += 1
                            
                    except Exception as e:
                        logger.warning(f"Failed to scrape {url}: {e}")
                        stats["failed"] += 1
                
                # Split content to chunks
                if content.strip():
                    chunks = text_splitter.split_text(content)
                    for i, chunk in enumerate(chunks):
                        results.append({
                            "content": chunk,
                            "metadata": {
                                "title": title,
                                "url": url,
                                "source": "web_search",
                                "chunk": i + 1,
                                "total_chunks": len(chunks),
                                "type": "web"
                            }
                        })
                    
                    stats["total_chunks"] += len(chunks)
                    
            except Exception as e:
                logger.error(f"Error processing search result {idx}: {e}")
                stats["failed"] += 1
            
    except Exception as e:
        logger.error(f"Web search error: {e}")
        st.error(f"❌ Error pencarian web: {str(e)}")
    
    st.session_state.processing_stats["web"] = stats
    return results

# ==========================
# 🧠 ChromaDB Vector Store Class
# ==========================
class ChromaVectorStore:
    """Enhanced ChromaDB vector store with caching and persistence"""
    
    def __init__(self, collection_name: str, embedding_model, chroma_client):
        self.client = chroma_client
        self.embedding_model = embedding_model
        self.collection_name = collection_name
        
        if not self.client or not self.embedding_model:
            raise ValueError("ChromaDB client and embedding model are required")
        
        # Reset collection for new data
        try:
            self.client.delete_collection(name=collection_name)
            logger.info(f"Deleted old collection: {collection_name}")
        except Exception:
            pass
        
        self.collection = self.client.get_or_create_collection(
            name=collection_name,
            metadata={"hnsw:space": "cosine"}
        )
        logger.info(f"Created collection: {collection_name}")
    
    def add_documents(self, docs: List[Dict[str, Any]]) -> int:
        """Add documents with batch processing"""
        if not docs:
            return 0
        
        batch_size = 100
        total_added = 0
        
        try:
            for i in range(0, len(docs), batch_size):
                batch = docs[i:i + batch_size]
                
                texts = [d["content"] for d in batch]
                metadatas = [d["metadata"] for d in batch]
                ids = [str(uuid.uuid4()) for _ in batch]
                
                # Generate embeddings
                embeddings = self.embedding_model.encode(
                    texts,
                    show_progress_bar=False,
                    convert_to_numpy=True
                ).tolist()
                
                # Add to collection
                self.collection.add(
                    ids=ids,
                    documents=texts,
                    metadatas=metadatas,
                    embeddings=embeddings
                )
                
                total_added += len(batch)
                logger.info(f"Added batch {i//batch_size + 1}: {len(batch)} docs")
            
            return total_added
            
        except Exception as e:
            logger.error(f"Failed to add documents: {e}")
            raise
    
    def retrieve(
        self,
        query: str,
        n_results: int = 5,
        filter_metadata: Optional[Dict] = None
    ) -> List[Dict[str, Any]]:
        """Retrieve documents with optional filtering"""
        try:
            query_embedding = self.embedding_model.encode(
                query,
                convert_to_numpy=True
            ).tolist()
            
            result = self.collection.query(
                query_embeddings=[query_embedding],
                n_results=min(n_results, self.collection.count()),
                include=["documents", "metadatas", "distances"],
                where=filter_metadata
            )
            
            retrieved_docs = []
            if result['ids'] and len(result['ids'][0]) > 0:
                for i in range(len(result['ids'][0])):
                    distance = result['distances'][0][i]
                    similarity = 1 / (1 + distance)
                    
                    retrieved_docs.append({
                        "content": result['documents'][0][i],
                        "metadata": result['metadatas'][0][i],
                        "score": similarity,
                        "distance": distance
                    })
            
            logger.info(f"Retrieved {len(retrieved_docs)} documents for query")
            return retrieved_docs
            
        except Exception as e:
            logger.error(f"Retrieval error: {e}")
            return []
    
    def get_stats(self) -> Dict[str, Any]:
        """Get collection statistics"""
        try:
            count = self.collection.count()
            return {
                "total_documents": count,
                "collection_name": self.collection_name
            }
        except Exception:
            return {"total_documents": 0, "collection_name": self.collection_name}

# ==========================
# 🤖 Gemini RAG System Class
# ==========================
class GeminiRAG:
    """Enhanced RAG system with conversation history"""
    
    def __init__(
        self,
        vector_store: ChromaVectorStore,
        api_key: str,
        model: str = "gemini-1.5-flash"
    ):
        self.vector_store = vector_store
        self.model = model
        self.api_key = api_key
    
    def _generate(self, prompt: str, temperature: float = 0.7) -> str:
        """Generate response from Gemini API"""
        if not self.api_key:
            return "❌ API Key tidak tersedia. Silakan masukkan API key di sidebar."
        
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{self.model}:generateContent?key={self.api_key}"
        headers = {"Content-Type": "application/json"}
        
        payload = {
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {
                "temperature": temperature,
                "maxOutputTokens": 2048,
                "topP": 0.95,
                "topK": 40
            },
            "safetySettings": [
                {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"}
            ]
        }
        
        try:
            response = requests.post(
                url,
                headers=headers,
                json=payload,
                timeout=60
            )
            response.raise_for_status()
            
            result = response.json()
            
            # Handle blocked content
            if "candidates" not in result or not result["candidates"]:
                return "❌ Respons diblokir oleh filter keamanan Gemini."
            
            candidate = result["candidates"][0]
            
            # Check if content was filtered
            if "content" not in candidate:
                finish_reason = candidate.get("finishReason", "UNKNOWN")
                return f"❌ Konten difilter. Alasan: {finish_reason}"
            
            answer = candidate["content"]["parts"][0]["text"]
            logger.info(f"Generated response: {len(answer)} chars")
            return answer
            
        except requests.exceptions.HTTPError as e:
            error_msg = f"API Error {e.response.status_code}"
            try:
                error_detail = e.response.json()
                error_msg += f": {error_detail.get('error', {}).get('message', 'Unknown error')}"
            except:
                pass
            logger.error(error_msg)
            return f"❌ {error_msg}"
            
        except Exception as e:
            logger.error(f"Generation error: {e}")
            return f"❌ Error: {str(e)}"
    
    def ask(
        self,
        query: str,
        n_results: int = 5,
        use_history: bool = True
    ) -> Dict[str, Any]:
        """Ask question with RAG"""
        
        # Retrieve relevant documents
        with st.spinner("🔍 Mencari informasi relevan..."):
            top_results = self.vector_store.retrieve(query, n_results=n_results)
        
        if not top_results:
            return {
                "answer": "Maaf, saya tidak menemukan informasi yang relevan untuk pertanyaan Anda. Coba tambahkan lebih banyak sumber data.",
                "contexts": [],
                "metadata": {"sources_used": 0}
            }
        
        # Build context
        context_parts = []
        for idx, r in enumerate(top_results, 1):
            source = r['metadata'].get('source', 'Unknown')
            title = r['metadata'].get('title', '')
            url = r['metadata'].get('url', '')
            
            source_info = f"{source}"
            if title:
                source_info = f"{title}"
            if url:
                source_info += f" ({url})"
            
            context_parts.append(
                f"[Sumber {idx}: {source_info}]\n{r['content']}"
            )
        
        context_text = "\n\n".join(context_parts)
        
        # Build conversation history context
        history_context = ""
        if use_history and st.session_state.chat_history:
            recent_history = st.session_state.chat_history[-3:]
            history_context = "\n\nRiwayat Percakapan:\n" + "\n".join([
                f"Q: {h['query']}\nA: {h['answer'][:200]}..."
                for h in recent_history
            ])
        
        # Generate prompt
        prompt = f"""Anda adalah asisten AI yang membantu menjawab pertanyaan berdasarkan konteks yang diberikan.

KONTEKS INFORMASI:
{context_text}
{history_context}

PERTANYAAN: {query}

INSTRUKSI:
- Jawab dalam bahasa Indonesia dengan jelas dan terstruktur
- Gunakan informasi dari konteks yang relevan
- Jika konteks tidak cukup untuk menjawab, jelaskan keterbatasannya
- Berikan jawaban yang informatif dan membantu
- Jangan menyebutkan "berdasarkan konteks" atau kalimat meta serupa

JAWABAN:"""
        
        # Generate response
        with st.spinner("🤖 Menghasilkan jawaban..."):
            answer = self._generate(prompt)
        
        # Save to history
        chat_entry = {
            "query": query,
            "answer": answer,
            "timestamp": datetime.now().isoformat(),
            "sources_used": len(top_results)
        }
        st.session_state.chat_history.append(chat_entry)
        
        return {
            "answer": answer,
            "contexts": top_results,
            "metadata": {
                "sources_used": len(top_results),
                "timestamp": chat_entry["timestamp"]
            }
        }

# ==========================
# 🎨 Streamlit UI
# ==========================
def main():
    st.title("🤖 Google Gemini RAG Assistant")
    st.caption("✨ Sistem RAG yang ditingkatkan dengan multiple data sources dan conversation memory")
    
    # Check dependencies
    if not check_dependencies():
        st.stop()
    
    # Initialize session state
    init_session_state()
    
    # Load resources
    EMBED_MODEL = load_embedding_model()
    CHROMA_CLIENT = load_chroma_client()
    TEXT_SPLITTER = get_text_splitter()
    
    if not EMBED_MODEL or not CHROMA_CLIENT:
        st.error("❌ Gagal load komponen utama sistem. Periksa instalasi dependencies.")
        st.stop()
    
    # Get API Key
    API_KEY = get_api_key()
    
    if not API_KEY:
        st.warning("⚠️ Silakan masukkan Gemini API Key di sidebar untuk melanjutkan.")
        st.info("💡 Dapatkan API key gratis di: https://ai.google.dev/")
        return
    
    # Sidebar Configuration
    st.sidebar.header("⚙️ Konfigurasi Sumber Data")
    
    # --- Document Upload ---
    use_documents = st.sidebar.checkbox("📄 Gunakan Dokumen", value=True)
    uploaded_files = None
    if use_documents:
        uploaded_files = st.sidebar.file_uploader(
            "📁 Upload dokumen",
            type=["pdf", "docx", "pptx", "txt", "md"],
            accept_multiple_files=True,
            key="doc_uploader",
            help="Support: PDF, DOCX, PPTX, TXT, MD"
        )
    
    # --- Web Search ---
    use_web_search = st.sidebar.checkbox("🌐 Gunakan Pencarian Web", value=False)
    web_query = ""
    scrape_full_page = False
    if use_web_search:
        web_query = st.sidebar.text_input(
            "🔎 Query pencarian web:",
            placeholder="Contoh: artificial intelligence trends 2024"
        )
        scrape_full_page = st.sidebar.checkbox(
            "🕷️ Scrape konten lengkap",
            value=True,
            help="Lebih lambat tapi lebih detail"
        )
    
    st.sidebar.markdown("---")
    
    # --- Advanced Settings ---
    with st.sidebar.expander("⚙️ Pengaturan Lanjutan"):
        n_results = st.slider(
            "Jumlah dokumen yang diambil",
            min_value=3,
            max_value=10,
            value=5,
            help="Lebih banyak = lebih komprehensif tapi lebih lambat"
        )
        
        use_history = st.checkbox(
            "Gunakan riwayat percakapan",
            value=True,
            help="AI akan mengingat 3 pertanyaan terakhir"
        )
    
    # --- Action Buttons ---
    col1, col2 = st.sidebar.columns(2)
    
    with col1:
        process_btn = st.button("🚀 Proses Data", type="primary", use_container_width=True)
    
    with col2:
        reset_btn = st.button("🔄 Reset", use_container_width=True)
    
    # --- Process Data ---
    if process_btn:
        all_docs = []
        st.session_state.processing_stats = {}
        
        # Process Documents
        if use_documents and uploaded_files:
            with st.spinner("📥 Memproses dokumen..."):
                temp_paths = []
                try:
                    for uf in uploaded_files:
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{uf.name}") as tmp:
                            tmp.write(uf.read())
                            temp_paths.append(tmp.name)
                    
                    document_chunks = load_and_split_documents(temp_paths, TEXT_SPLITTER)
                    all_docs.extend(document_chunks)
                    
                    if document_chunks:
                        st.success(
                            f"✅ Berhasil memproses {len(uploaded_files)} dokumen → "
                            f"{len(document_chunks)} chunks"
                        )
                finally:
                    # Cleanup temp files
                    for p in temp_paths:
                        try:
                            os.remove(p)
                        except:
                            pass
        
        # Process Web Search
        if use_web_search and web_query.strip():
            with st.spinner("🌐 Mencari di web..."):
                web_chunks = search_and_scrape_web(
                    web_query,
                    TEXT_SPLITTER,
                    max_results=3,
                    scrape_full_page=scrape_full_page
                )
                all_docs.extend(web_chunks)
                
                if web_chunks:
                    st.success(f"✅ Berhasil mengambil {len(web_chunks)} chunks dari web")
        
        # Initialize Vector Store
        if all_docs:
            with st.spinner("🧠 Menyimpan ke vector database..."):
                try:
                    vector_store = ChromaVectorStore(
                        collection_name="rag_collection",
                        embedding_model=EMBED_MODEL,
                        chroma_client=CHROMA_CLIENT
                    )
                    added = vector_store.add_documents(all_docs)
                    
                    st.session_state.vector_store = vector_store
                    st.session_state.initialized = True
                    
                    # Show stats
                    stats = vector_store.get_stats()
                    st.success(
                        f"🎉 Sistem RAG siap! Total dokumen: {stats['total_documents']}"
                    )
